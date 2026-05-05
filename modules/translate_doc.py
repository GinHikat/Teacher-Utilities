import os, sys
from pathlib import Path
from docx import Document
from pptx import Presentation
from deep_translator import GoogleTranslator
from tqdm import tqdm
import win32com.client

project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
if project_root not in sys.path:
    sys.path.append(project_root)

def translate_text(text, target_lang="en"):
    """Translate text to target language."""
    if not text or not text.strip():
        return text or ""

    try:
        translated = GoogleTranslator(source="vi", target=target_lang).translate(text)
        return translated if translated is not None else text
    except Exception as e:
        print(f"Translation error: {text[:30]}... -> {e}")
        return text

def convert_doc_to_docx(input_path: Path) -> Path:
    """Convert .doc file to .docx using Microsoft Word."""
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False

    doc = word.Documents.Open(str(input_path))
    output_path = input_path.with_suffix(".docx")

    doc.SaveAs(str(output_path), FileFormat=16)  # 16 = docx
    doc.Close()
    word.Quit()

    return output_path

def translate_paragraph(paragraph):
    """Translate a single paragraph while preserving formatting."""
    text = paragraph.text.strip()
    if not text:
        return

    translated = translate_text(text)
    if translated is None:
        translated = text

    if not paragraph.runs:
        if hasattr(paragraph, "add_run"):
            # docx style
            try:
                paragraph.add_run(translated)
            except TypeError:
                # pptx style
                paragraph.add_run().text = translated
        return

    # Clear existing runs but keep structure
    for run in paragraph.runs:
        run.text = ""

    paragraph.runs[0].text = translated

def translate_full_presentation(input_path: Path, output_path: Path):
    """Translate an entire PowerPoint presentation."""
    if not input_path.exists():
        print(f"File not found: {input_path}")
        return

    print(f"\nOpening presentation: {input_path}")

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"Cannot open presentation: {e}")
        return

    # Iterate through slides
    print(f"Translating {len(prs.slides)} slides...")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    translate_paragraph(paragraph)
            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            translate_paragraph(paragraph)

    prs.save(str(output_path))
    print(f"Saved translated file: {output_path}")

def translate_full_document(input_path: Path, output_path: Path):
    """Translate an entire Word document."""
    if not input_path.exists():
        print(f"File not found: {input_path}")
        return

    print(f"\nOpening document: {input_path}")

    try:
        doc = Document(str(input_path))
    except Exception as e:
        print(f"Cannot open document: {e}")
        return

    # Translate body paragraphs
    print(f"Translating {len(doc.paragraphs)} paragraphs...")
    for paragraph in doc.paragraphs:
        translate_paragraph(paragraph)

    # Translate tables
    print(f"Translating {len(doc.tables)} tables...")
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    translate_paragraph(paragraph)

    doc.save(str(output_path))
    print(f"Saved translated file: {output_path}")

def main():

    input_folder = Path(
        os.path.join(os.path.dirname(__file__), "..", "files")
    )

    target_folder = input_folder / "trans"
    target_folder.mkdir(exist_ok=True)

    # Find already translated files
    translated_files = {
        p.stem.replace("_trans", "") 
        for p in target_folder.glob("*_trans.*")
    }

    # Collect input documents
    files = [
        f for f in input_folder.iterdir()
        if f.is_file() and f.suffix.lower() in [".doc", ".docx", ".pptx", ".pptm"]
    ]

    for file in tqdm(files, desc="Processing files"):

        if file.stem in translated_files:
            continue

        input_path = file

        # Convert .doc to .docx if needed
        if file.suffix.lower() == ".doc":
            print(f"Converting {file.name} -> docx")
            input_path = convert_doc_to_docx(file)

        if input_path.suffix.lower() == ".docx":
            output_file = target_folder / f"{input_path.stem}_trans.docx"
            translate_full_document(input_path, output_file)
        elif input_path.suffix.lower() in [".pptx", ".pptm"]:
            suffix = input_path.suffix.lower()
            output_file = target_folder / f"{input_path.stem}_trans{suffix}"
            translate_full_presentation(input_path, output_file)


if __name__ == "__main__":
    main()